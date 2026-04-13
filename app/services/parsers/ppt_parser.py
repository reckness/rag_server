"""
PPT/PPTX 文档解析器
每张幻灯片作为一个节点
"""
import os
from typing import Dict, Any, List
from .base_parser import BaseDocumentParser


class PptParser(BaseDocumentParser):
    """PPT/PPTX 文档解析器"""

    def parse(self, file_path: str, doc_name: str = None) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        if doc_name is None:
            doc_name = os.path.splitext(os.path.basename(file_path))[0]
        prs = Presentation(file_path)

        structure = []
        slide_counter = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_counter += 1
            slide_title = ""
            slide_texts: List[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # 提取标题（占位符类型 = 标题）
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue

                text_frame = shape.text_frame
                shape_text = ""
                for para in text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        shape_text += para_text + "\n"

                shape_text = shape_text.strip()
                if not shape_text:
                    continue

                # 判断是否为标题占位符
                ph_idx = getattr(getattr(shape, "placeholder_format", None), "idx", None)
                if ph_idx == 0 and not slide_title:
                    slide_title = shape_text
                else:
                    slide_texts.append(shape_text)

            if not slide_title:
                slide_title = f"幻灯片 {slide_num}"

            slide_content = "\n".join(slide_texts)

            node = {
                "title": slide_title,
                "node_id": str(slide_counter).zfill(4),
                "text": slide_content,
                "start_index": slide_num,
                "end_index": slide_num,
            }
            structure.append(node)

        if not structure:
            structure = [{
                "title": doc_name,
                "node_id": "0001",
                "text": "",
                "start_index": 1,
                "end_index": 1,
            }]

        return {
            "doc_name": doc_name,
            "structure": structure,
        }
